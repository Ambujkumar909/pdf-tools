import os
from flask import Flask, request, redirect, url_for, render_template, send_from_directory, flash
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from io import BytesIO
from docx import Document
import time

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'
app.config['MAX_CONTENT_LENGTH'] = 5 * 1024 * 1024 * 1024  # 5 GB limit
app.config['APP_VERSION'] = '1.0.0'

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
app.secret_key = '46c7e0c9899d5b5485b4c7e72da034f7'  # Replace with your own secret key

def split_pdf(input_pdf_path, output_folder):
    pdf_document = fitz.open(input_pdf_path)
    output_files = []

    for page_num in range(len(pdf_document)):
        output_pdf_path = os.path.join(output_folder, f"page_{page_num + 1}.pdf")
        new_document = fitz.open()
        new_document.insert_pdf(pdf_document, from_page=page_num, to_page=page_num)
        new_document.save(output_pdf_path)
        new_document.close()
        output_files.append(output_pdf_path)

    pdf_document.close()
    return output_files

def merge_pdfs(pdf_list, output_path):
    merged_document = fitz.open()
    for pdf in pdf_list:
        with fitz.open(pdf) as document:
            merged_document.insert_pdf(document)
    merged_document.save(output_path)

def pdf_to_excel(input_pdf, output_folder):
    output_excel = os.path.join(output_folder, 'converted.xlsx')
    pdf_document = fitz.open(input_pdf)
    with pd.ExcelWriter(output_excel, engine='xlsxwriter') as writer:
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text = page.get_text()
            paragraphs = text.split('\n\n')
            df = pd.DataFrame({"Content": paragraphs})
            sheet_name = f'Page_{page_num + 1}'
            df.to_excel(writer, sheet_name=sheet_name, index=False)
    return output_excel

def pdf_to_images(input_pdf, output_folder):
    output_images_folder = os.path.join(output_folder, 'images')
    os.makedirs(output_images_folder, exist_ok=True)
    images = []
    with fitz.open(input_pdf) as pdf_document:
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            timestamp = int(time.time() * 1000)
            pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))  # Increase DPI for better quality
            image_path = os.path.join(output_images_folder, f"page_{page_num + 1}_{timestamp}.jpg")
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img.save(image_path, "JPEG", quality=100, dpi=(300, 300))  # Set DPI and quality
            images.append(image_path)
    return output_images_folder, images

def pdf_to_ppt(input_pdf, output_folder):
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    output_pptx = os.path.join(output_folder, 'converted.pptx')

    with fitz.open(input_pdf) as pdf_document:
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))  # Increase DPI for better quality
            image_bytes = pix.tobytes()
            image = Image.open(BytesIO(image_bytes))
            temp_image_path = f"temp_image_{page_num + 1}.png"
            image.save(temp_image_path, "PNG", dpi=(300, 300))  # Set DPI for image quality

            slide_layout = prs.slide_layouts[6]  # Blank slide layout (custom)
            slide = prs.slides.add_slide(slide_layout)

            img_left = Inches(0.5)
            img_top = Inches(0.5)
            img_width = slide_width - Inches(1)
            img_height = slide_height - Inches(1)
            slide.shapes.add_picture(temp_image_path, img_left, img_top, width=img_width, height=img_height)

            os.remove(temp_image_path)  # Remove temporary image file

    prs.save(output_pptx)
    return output_pptx

def pdf_to_word(input_pdf, output_folder):
    output_docx = os.path.join(output_folder, 'converted.docx')
    doc = Document()
    with fitz.open(input_pdf) as pdf_document:
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))  # Increase DPI for better quality
            image_bytes = pix.tobytes()
            image = Image.open(BytesIO(image_bytes))
            temp_image_path = f"temp_image_{page_num + 1}.png"
            image.save(temp_image_path, "PNG", dpi=(150, 150))  # Set DPI for image quality
            doc.add_picture(temp_image_path, width=Inches(6))
            os.remove(temp_image_path)
    doc.save(output_docx)
    return output_docx

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() == 'pdf'

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/split', methods=['GET', 'POST'])
def split_pdf_route():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file part')
            return redirect(request.url)
        file = request.files['file']
        if file.filename == '':
            flash('No selected file')
            return redirect(request.url)
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            output_files = split_pdf(file_path, app.config['OUTPUT_FOLDER'])
            return render_template('result.html', output_files=[os.path.basename(f) for f in output_files], operation='split')
        else:
            flash('Invalid file type. Only PDF files are allowed.')
            return redirect(request.url)
    return render_template('split.html')

@app.route('/merge', methods=['GET', 'POST'])
def merge_pdf_route():
    if request.method == 'POST':
        files = request.files.getlist('files[]')
        if len(files) < 2:
            flash('Select at least two PDF files to merge')
            return redirect(request.url)
        pdf_list = []
        for file in files:
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                pdf_list.append(filepath)
        if pdf_list:
            merged_pdf_path = os.path.join(app.config['OUTPUT_FOLDER'], 'merged.pdf')
            merge_pdfs(pdf_list, merged_pdf_path)
            return render_template('result.html', output_file='merged.pdf', operation='merge')
        else:
            flash('Invalid file type. Only PDF files are allowed.')
            return redirect(request.url)
    return render_template('merge.html')

@app.route('/convert', methods=['GET', 'POST'])
def convert_pdf():
    if request.method == 'POST':
        try:
            conversion_choice = request.form['conversion_choice']
            file = request.files['file']

            if file.filename == '':
                flash('No selected file')
                return redirect(request.url)

            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(file_path)

                output_folder = app.config['OUTPUT_FOLDER']

                if conversion_choice == 'excel':
                    output_excel = pdf_to_excel(file_path, output_folder)
                    if output_excel:
                        return render_template('result.html', output_file=os.path.basename(output_excel), operation='convert')
                    else:
                        flash('Conversion to Excel failed')
                        return redirect(request.url)

                elif conversion_choice == 'images':
                    output_images_folder, images = pdf_to_images(file_path, output_folder)
                    if output_images_folder:
                        return render_template('result.html', output_files=[os.path.basename(img) for img in images],
                                               operation='convert')
                    else:
                        flash('Conversion to images failed')
                        return redirect(request.url)


                elif conversion_choice == 'ppt':
                    output_pptx = pdf_to_ppt(file_path, output_folder)
                    if output_pptx:
                        return render_template('result.html', output_file=os.path.basename(output_pptx), operation='convert')
                    else:
                        flash('Conversion to PowerPoint failed')
                        return redirect(request.url)

                elif conversion_choice == 'word':
                    output_docx = pdf_to_word(file_path, output_folder)
                    if output_docx:
                        return render_template('result.html', output_file=os.path.basename(output_docx), operation='convert')
                    else:
                        flash('Conversion to Word failed')
                        return redirect(request.url)

                else:
                    flash('Invalid conversion choice')
                    return redirect(request.url)

            else:
                flash('Invalid file type. Only PDF files are allowed.')
                return redirect(request.url)

        except Exception as e:
            flash(f'An error occurred during conversion: {str(e)}')
            return redirect(request.url)

    return render_template('convert.html')

@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory(app.config['OUTPUT_FOLDER'], filename)

if __name__ == "__main__":
    app.run(debug=True)
